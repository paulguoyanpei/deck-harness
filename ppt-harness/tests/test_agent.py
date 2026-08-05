from __future__ import annotations

import io
import itertools
import json
import os
import shutil
import sys
import tempfile
import unittest
from copy import deepcopy
from pathlib import Path
from unittest.mock import patch


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.append(str(ROOT / "pptx" / "scripts"))
sys.path.append(str(ROOT / "pptx" / "scripts" / "office"))

import layout_check  # noqa: E402
import ppt_agent  # noqa: E402
from validators.pptx import PPTXSchemaValidator  # noqa: E402


class FakeMessages:
    def __init__(self, responses):
        self.responses = iter(responses)
        self.calls = []

    def create(self, **kwargs):
        self.calls.append(deepcopy(kwargs))
        return next(self.responses)


class FakeClient:
    def __init__(self, responses):
        self.messages = FakeMessages(responses)


class ToolRunnerTests(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.runner = ppt_agent.ToolRunner(self.root)

    def tearDown(self):
        self.temp.cleanup()

    def test_write_read_and_unique_edit(self):
        result, error = self.runner.execute(
            "write_file", {"path": "src/deck.py", "content": "alpha\nbeta\n"}
        )
        self.assertFalse(error)
        self.assertIn("File created successfully at:", result)
        self.assertIn("no need to read it back", result)

        result, error = self.runner.execute(
            "read_file", {"path": "src/deck.py", "offset": 1, "limit": 1}
        )
        self.assertFalse(error)
        self.assertIn("     2  beta", result)

        result, error = self.runner.execute(
            "edit_file",
            {"path": "src/deck.py", "old_text": "beta", "new_text": "gamma"},
        )
        self.assertFalse(error)
        self.assertEqual((self.root / "src/deck.py").read_text(), "alpha\ngamma\n")
        self.assertIn("File updated successfully at:", result)

    def test_edit_rejects_ambiguous_match(self):
        (self.root / "deck.py").write_text("same same", encoding="utf-8")
        result, error = self.runner.execute(
            "edit_file",
            {"path": "deck.py", "old_text": "same", "new_text": "new"},
        )
        self.assertTrue(error)
        self.assertIn("matched 2 times", result)

    def test_file_tools_cannot_escape_workdir(self):
        result, error = self.runner.execute(
            "write_file", {"path": "../escape.txt", "content": "no"}
        )
        self.assertTrue(error)
        self.assertIn("inside the working directory", result)

    def test_bash_strips_api_keys(self):
        with patch.dict(os.environ, {"DEEPSEEK_API_KEY": "super-secret"}):
            result = self.runner.bash('test -z "$DEEPSEEK_API_KEY"')
        self.assertEqual(result.exit_code, 0)

    def test_bash_reports_failure_and_timeout(self):
        failed = self.runner.bash("echo bad >&2; exit 7")
        self.assertEqual(failed.exit_code, 7)
        self.assertIn("bad", failed.stderr)

        timed_out = self.runner.bash("sleep 2", timeout_seconds=1)
        self.assertTrue(timed_out.timed_out)
        self.assertEqual(timed_out.exit_code, 124)

    def test_bash_model_content_is_concise_but_logs_are_structured(self):
        execution = self.runner.execute("bash", {"command": "printf out; printf err >&2; exit 7"})
        self.assertTrue(execution.is_error)
        self.assertEqual(execution.model_content, "Exit code 7\nout\nerr")
        self.assertNotIn('"command"', execution.model_content)
        self.assertEqual(execution.details["command"], "printf out; printf err >&2; exit 7")
        self.assertEqual(execution.details["exit_code"], 7)
        self.assertEqual(execution.details["stdout"], "out")
        self.assertEqual(execution.details["stderr"], "err")

        empty_success = self.runner.execute("bash", {"command": "true"})
        self.assertEqual(empty_success.model_content, "Command completed successfully.")

    def test_tool_output_is_truncated(self):
        result = self.runner.bash("python3 -c 'print(\"x\" * 40000)'")
        self.assertLessEqual(len(result.stdout), ppt_agent.MAX_TOOL_OUTPUT + 100)
        self.assertIn("chars omitted", result.stdout)


class OfficeValidatorTests(unittest.TestCase):
    def test_duplicate_cnvpr_ids_are_rejected(self):
        with tempfile.TemporaryDirectory() as temp:
            slide = Path(temp) / "ppt" / "slides" / "slide1.xml"
            slide.parent.mkdir(parents=True)
            slide.write_text(
                '<p:sld xmlns:p="urn:p"><p:cNvPr id="11" name="table"/>'
                '<p:cNvPr id="11" name="icon"/></p:sld>',
                encoding="utf-8",
            )
            with patch("sys.stdout", new=io.StringIO()) as out:
                valid = PPTXSchemaValidator(temp).validate_unique_ids()

        self.assertFalse(valid)
        self.assertIn("Duplicate id='11' in <cnvpr>", out.getvalue())


class AgentLoopTests(unittest.TestCase):
    def test_tool_results_and_failed_qa_return_to_model(self):
        responses = [
            {
                "stop_reason": "tool_use",
                "content": [
                    {
                        "type": "tool_use",
                        "id": "tool-1",
                        "name": "write_file",
                        "input": {"path": "deck.py", "content": "print('one')\n"},
                    }
                ],
            },
            {"stop_reason": "end_turn", "content": [{"type": "text", "text": "Done"}]},
            {
                "stop_reason": "tool_use",
                "content": [
                    {
                        "type": "tool_use",
                        "id": "tool-2",
                        "name": "edit_file",
                        "input": {
                            "path": "deck.py",
                            "old_text": "one",
                            "new_text": "two",
                        },
                    }
                ],
            },
            {"stop_reason": "end_turn", "content": [{"type": "text", "text": "Fixed"}]},
        ]
        client = FakeClient(responses)
        with tempfile.TemporaryDirectory() as temp, patch.object(
            ppt_agent, "qa_presentation", side_effect=[(False, "FAIL\nmissing"), (True, "PASS\n")]
        ):
            final = ppt_agent.run_agent(client, "Build slides", Path(temp), max_steps=6)
            log_lines = (Path(temp) / "run.jsonl").read_text().splitlines()

        self.assertEqual(final, "Fixed")
        self.assertEqual(len(client.messages.calls), 4)
        official_skill = ppt_agent.DEFAULT_SKILL_PATH.read_text(encoding="utf-8")
        self.assertIn(official_skill, client.messages.calls[0]["system"])
        self.assertIn(str(ppt_agent.DEFAULT_SKILL_PATH.parent), client.messages.calls[0]["system"])
        self.assertIn(str(Path(temp).resolve()), client.messages.calls[0]["system"])
        self.assertNotIn("HARNESS_PROMPT", client.messages.calls[0]["system"])
        self.assertIn("use programmatic QA where needed", client.messages.calls[0]["system"])
        self.assertNotIn("Do not run validate.py", client.messages.calls[0]["system"])
        second_messages = client.messages.calls[1]["messages"]
        self.assertEqual(second_messages[-1]["content"][0]["tool_use_id"], "tool-1")
        self.assertIn("File created successfully at:", second_messages[-1]["content"][0]["content"])
        tool_log = next(json.loads(line) for line in log_lines if json.loads(line)["kind"] == "tool")
        self.assertIn("model_content", tool_log["data"])
        self.assertIn("details", tool_log["data"])
        third_messages = client.messages.calls[2]["messages"]
        self.assertIn("Deterministic QA failed", third_messages[-1]["content"])
        self.assertTrue(any(json.loads(line)["kind"] == "qa" for line in log_lines))

    def test_missing_skill_fails_before_api_call(self):
        client = FakeClient([])
        with tempfile.TemporaryDirectory() as temp:
            missing = Path(temp) / "missing-skill.md"
            with self.assertRaisesRegex(RuntimeError, "PPTX skill not found"):
                ppt_agent.run_agent(client, "Build", Path(temp) / "run", skill_path=missing)
        self.assertEqual(client.messages.calls, [])

    def test_max_steps_is_enforced(self):
        response = {
            "stop_reason": "tool_use",
            "content": [
                {
                    "type": "tool_use",
                    "id": "tool-1",
                    "name": "bash",
                    "input": {"command": "true"},
                }
            ],
        }
        client = FakeClient([response])
        with tempfile.TemporaryDirectory() as temp:
            with self.assertRaisesRegex(RuntimeError, "maximum of 1 steps"):
                ppt_agent.run_agent(client, "Build", Path(temp), max_steps=1)

    def test_relative_workdir_is_resolved_before_qa(self):
        client = FakeClient(
            [{"stop_reason": "end_turn", "content": [{"type": "text", "text": "Done"}]}]
        )
        with tempfile.TemporaryDirectory() as temp:
            parent = Path(temp)
            old_cwd = Path.cwd()
            try:
                os.chdir(parent)
                with patch.object(ppt_agent, "qa_presentation", return_value=(True, "PASS\n")) as qa:
                    ppt_agent.run_agent(client, "Build", Path("runs/demo"), max_steps=1)
            finally:
                os.chdir(old_cwd)
            qa_workdir = qa.call_args.args[0]
            self.assertTrue(qa_workdir.is_absolute())
            self.assertEqual(qa_workdir, (parent / "runs/demo").resolve())

    def test_default_step_limit_is_40(self):
        args = ppt_agent.parse_args(["--out", "runs/demo", "Build"])
        self.assertEqual(args.max_steps, 40)


class ContentBalanceTests(unittest.TestCase):
    def test_balanced_rectangles_pass_metric(self):
        metrics = layout_check.content_balance_metrics(
            [(0.6, 1.0, 5.5, 6.5), (7.8, 1.0, 12.7, 6.5)], 13.333, 7.5
        )
        self.assertGreater(metrics["left"], 25)
        self.assertGreater(metrics["right"], 25)
        self.assertLess(metrics["empty_vertical"], 0.25)

    def test_left_heavy_rectangles_expose_empty_band(self):
        metrics = layout_check.content_balance_metrics([(0.6, 1.0, 5.8, 6.5)], 13.333, 7.5)
        self.assertGreater(metrics["left"], 25)
        self.assertLess(metrics["right"], 5)
        self.assertGreaterEqual(metrics["empty_vertical"], 0.25)

    def test_tight_text_bounds_detect_severe_one_sided_slide(self):
        from pptx import Presentation
        from pptx.util import Inches, Pt

        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            deck = Presentation()
            deck.slide_width, deck.slide_height = Inches(13.333), Inches(7.5)
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            for row in range(9):
                box = slide.shapes.add_textbox(
                    Inches(0.65), Inches(1.35 + row * 0.5), Inches(11.4), Inches(0.35)
                )
                box.text = f"{row + 1}. Verify facts against a trusted source."
                box.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
            deck.save(root / "deck.pptx")

            passed, detail = layout_check.analyse_balance(root / "deck.pptx")

            self.assertFalse(passed)
            self.assertIn("FAIL  slide 1: right body", detail)
            self.assertIn("Redistribute or resize existing content", detail)


@unittest.skipUnless(
    all(shutil.which(name) for name in ["soffice", "pdfinfo", "pdftotext", "pdftoppm", "montage", "unzip"]),
    "presentation QA system tools are not installed",
)
class QaIntegrationTests(unittest.TestCase):
    def _make_deck(self, root: Path, text: str = "Hello") -> None:
        from pptx import Presentation

        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        box = slide.shapes.add_textbox(1_000_000, 1_000_000, 5_000_000, 1_000_000)
        box.text = text
        deck.save(root / "deck.pptx")

    def test_valid_deck_creates_all_artifacts(self):
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            self._make_deck(root)
            old_cwd = Path.cwd()
            try:
                os.chdir(root.parent)
                passed, report = ppt_agent.qa_presentation(Path(root.name))
            finally:
                os.chdir(old_cwd)
            self.assertTrue(passed, report)
            self.assertNotIn("generator", report)
            self.assertTrue((root / "deck.pdf").is_file())
            self.assertTrue((root / "preview.jpg").is_file())
            self.assertEqual(len(list((root / "slides").glob("*.jpg"))), 1)

    def test_placeholder_fails_qa(self):
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            self._make_deck(root, "TODO replace me")
            (root / "deck.pdf").write_bytes(b"old pdf")
            (root / "preview.jpg").write_bytes(b"old preview")
            (root / "slides").mkdir()
            (root / "slides/old.jpg").write_bytes(b"old slide")
            passed, report = ppt_agent.qa_presentation(root)
            self.assertFalse(passed)
            self.assertIn("[FAIL] placeholders", report)
            self.assertEqual((root / "deck.pdf").read_bytes(), b"old pdf")
            self.assertEqual((root / "preview.jpg").read_bytes(), b"old preview")
            self.assertEqual((root / "slides/old.jpg").read_bytes(), b"old slide")

    def test_broken_pptx_fails_qa(self):
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            (root / "deck.pptx").write_text("not a zip")
            passed, report = ppt_agent.qa_presentation(root)
            self.assertFalse(passed)
            self.assertIn("[FAIL] pptx_zip", report)


class LayoutCheckTests(unittest.TestCase):
    """pptx/scripts/layout_check.py — geometry and font-metric layout checks."""

    @staticmethod
    def _deck(root: Path, build) -> Path:
        from pptx import Presentation
        from pptx.util import Inches

        deck = Presentation()
        deck.slide_width, deck.slide_height = Inches(10), Inches(5.625)
        build(deck.slides.add_slide(deck.slide_layouts[6]))
        path = root / "deck.pptx"
        deck.save(path)
        return path

    @staticmethod
    def _shape(slide, name, x, y, w, h):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.name = name
        return shape

    @staticmethod
    def _text(slide, name, x, y, w, h, text, size=12, font=None, wrap=False):
        from pptx.util import Inches, Pt

        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.name = name
        if wrap:
            box.text_frame.word_wrap = True
            box.text_frame.auto_size = None
        box.text_frame.text = text
        for run in box.text_frame.paragraphs[0].runs:
            run.font.size = Pt(size)
            if font:
                run.font.name = font
        return box

    def _report(self, build) -> str:
        with tempfile.TemporaryDirectory() as temp:
            return layout_check.analyse(str(self._deck(Path(temp), build)))

    # -- tree building ----------------------------------------------------

    def test_background_prefix_is_never_reported(self):
        def build(slide):
            self._shape(slide, "@bg/wash", 0, 0, 10, 5.625)
            self._shape(slide, "card", 1, 1, 3, 2)

        self.assertIn("0 warnings", self._report(build))

    def test_declared_child_on_its_parent_is_not_an_overlap(self):
        def build(slide):
            self._shape(slide, "card", 0.5, 0.5, 4, 2)
            self._text(slide, "card/title", 0.7, 0.7, 3, 0.4, "Title")

        self.assertIn("0 warnings", self._report(build))

    def test_identical_rectangles_nest_by_z_order(self):
        # A number drawn on top of its circle shares the circle's rectangle.
        def build(slide):
            self._shape(slide, "Shape 0", 1, 1, 0.44, 0.44)
            self._text(slide, "Text 1", 1, 1, 0.44, 0.44, "1")

        report = self._report(build)
        self.assertIn("0 warnings", report)
        self.assertIn("named 0 | inferred 2", report)

    def test_unnamed_element_is_adopted_by_smallest_named_container(self):
        def build(slide):
            self._shape(slide, "outer", 0, 0, 9, 5)
            self._shape(slide, "card", 1, 1, 4, 2)
            self._text(slide, "Text 7", 1.2, 1.2, 1, 0.3, "hi")

        self.assertIn("0 warnings", self._report(build))

    # -- overlap ----------------------------------------------------------

    def test_partial_overlap_of_siblings_is_reported(self):
        # The slide-3 defect: dots run to 6.25, the pill starts at 6.15.
        def build(slide):
            self._shape(slide, "row/dot", 6.11, 3.24, 0.14, 0.14)
            self._shape(slide, "row/pill", 6.15, 3.16, 3.08, 0.42)

        report = self._report(build)
        self.assertIn("row/dot × row/pill", report)
        self.assertIn("overlapping, move 0.10in to clear", report)

    def test_severity_is_how_far_something_must_move(self):
        # A hairline overlap is a nudge; a buried element is severe. Coverage
        # ratio cannot tell those apart -- it calls both of these ~100%.
        def build(slide):
            self._shape(slide, "a", 1.0, 1.0, 2.0, 1.0)
            self._shape(slide, "a2", 2.97, 1.0, 2.0, 1.0)    # 0.03in overlap
            self._shape(slide, "b", 1.0, 3.0, 2.0, 1.0)
            self._shape(slide, "b2", 2.0, 3.0, 2.0, 1.0)     # 1.00in overlap

        report = self._report(build)
        self.assertIn("move 0.03in to clear", report)
        self.assertNotIn("move 0.03in to clear  [severe]", report)
        self.assertIn("move 1.00in to clear  [severe]", report)

    def test_abutting_edges_must_still_stand_clear(self):
        # Under the binary rule a 0.02in sliver is no longer waved through as
        # "merely abutting" -- elements either nest or keep their distance.
        def build(slide):
            self._shape(slide, "a", 1, 1, 2, 1.0)
            self._shape(slide, "b", 1, 1.98, 2, 1.0)   # 0.02in sliver

        self.assertIn("WARN  overlap", self._report(build))

    def test_elements_merely_close_are_reported_as_clearance(self):
        def build(slide):
            self._shape(slide, "a", 1, 1, 2, 1.0)
            self._shape(slide, "b", 3.03, 1, 2, 1.0)   # 0.03in apart

        report = self._report(build)
        self.assertIn("WARN  clearance", report)
        self.assertIn("0.03in apart, needs 0.05in", report)

    def test_elements_a_clear_gap_apart_pass(self):
        def build(slide):
            self._shape(slide, "a", 1, 1, 2, 1.0)
            self._shape(slide, "b", 3.08, 1, 2, 1.0)   # 0.08in apart

        self.assertIn("0 warnings", self._report(build))

    def test_one_visual_object_is_reported_once(self):
        # A pill and the text drawn on it must not both be reported against a bar.
        def build(slide):
            self._shape(slide, "bar", 1.0, 1.0, 0.5, 2.0)
            self._shape(slide, "callout", 1.2, 1.0, 2.0, 0.5)
            self._text(slide, "callout/label", 1.2, 1.0, 2.0, 0.5, "note")

        self.assertEqual(self._report(build).count("WARN  overlap"), 1)

    # -- escape, fit, center, baseline ------------------------------------

    def test_text_escaping_its_parent_is_an_overlap(self):
        # There is no separate escape check: hanging off a parent's edge is
        # simply a pair that neither nests nor stands clear.
        def build(slide):
            self._shape(slide, "card", 0.5, 0.5, 2, 1)
            self._text(slide, "card/wide", 0.7, 0.7, 3, 0.4,
                       "text far too long to stay inside")

        report = self._report(build)
        self.assertIn("WARN  overlap", report)
        self.assertIn("card", report)

    def test_a_loose_frame_hanging_out_is_not_an_escape(self):
        # The frame runs 1.2in past the card, but the glyphs stop well inside
        # it. Judging the frame reported this; judging the ink does not.
        def build(slide):
            self._shape(slide, "card", 0.5, 0.5, 2, 1)
            self._text(slide, "card/short", 0.7, 0.7, 3, 0.4, "short")

        self.assertIn("0 warnings", self._report(build))

    def test_text_too_wide_for_its_box_is_reported(self):
        # The page-number defect: 0.5in box minus 0.2in of insets holds 0.30in.
        def build(slide):
            self._text(slide, "footer/pageno", 8.95, 5.3, 0.5, 0.22,
                       "02 / 10", size=9, font="Microsoft YaHei", wrap=True)

        report = self._report(build)
        self.assertIn("footer/pageno", report)
        self.assertIn("needs 2 lines, box holds 1", report)

    def test_text_that_fits_is_not_reported(self):
        def build(slide):
            self._text(slide, "footer/pageno", 8.0, 5.3, 1.5, 0.3,
                       "02 / 10", size=9, font="Microsoft YaHei", wrap=True)

        self.assertIn("0 warnings", self._report(build))

    def test_tall_top_anchored_text_box_is_reported_as_underfilled(self):
        def build(slide):
            self._text(slide, "card/desc", 1, 1, 4, 1.4,
                       "A short body paragraph leaves most of this tall frame empty.",
                       size=12, wrap=True)

        report = self._report(build)
        self.assertIn("WARN  underfill", report)
        self.assertIn("expected at least 40%", report)
        self.assertIn("Increase the font size or add more text content", report)
        self.assertIn("peer boxes in the same group", report)

    def test_middle_anchored_banner_is_not_reported_as_underfilled(self):
        from pptx.enum.text import MSO_ANCHOR

        def build(slide):
            box = self._text(slide, "takeaway/text", 1, 1, 8, 1.4,
                             "A deliberately centred takeaway banner has spare space.",
                             size=12, wrap=True)
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        self.assertNotIn("WARN  underfill", self._report(build))

    def test_centred_text_is_not_measured_against_the_backdrop(self):
        # Parents are geometric now, so anything not nested in a card lands on
        # the full-slide background. A bottom caption is horizontally centred
        # and vertically nowhere near the slide's middle -- by design.
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

        def build(slide):
            self._shape(slide, "@bg/slide-bg", 0, 0, 10, 5.625)
            box = self._text(slide, "caption", 1.0, 4.8, 8.0, 0.4, "a bottom note")
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.assertNotIn("WARN  center", self._report(build))

    def test_centred_text_off_its_shape_centre_is_reported(self):
        # The badge defect: circle at y=0.64, its label at y=0.62.
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

        def build(slide):
            self._shape(slide, "badge", 8.85, 0.64, 0.36, 0.36)
            box = self._text(slide, "badge/mark", 8.85, 0.62, 0.36, 0.36, "?")
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        report = self._report(build)
        self.assertIn("badge/mark", report)
        self.assertIn("off the centre of badge", report)

    def test_row_with_mismatched_text_anchors_is_reported(self):
        # The slide-8 defect: a top-anchored title beside a middle-anchored body.
        from pptx.enum.text import MSO_ANCHOR

        def build(slide):
            self._shape(slide, "row", 0.55, 1.5, 8.9, 0.58)
            self._text(slide, "row/title", 1.42, 1.53, 1.8, 0.32, "Title", size=14)
            body = self._text(slide, "row/body", 3.3, 1.53, 5.9, 0.5, "Body", size=12)
            body.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        report = self._report(build)
        self.assertIn("WARN  baseline", report)
        self.assertIn("share a row", report)

    # -- reporting --------------------------------------------------------

    def _overlapping_pair(self, slide):
        self._shape(slide, "chart/bar", 5.0, 0.5, 1.0, 2.0)
        self._shape(slide, "chart/callout", 5.5, 0.5, 2.0, 0.5)

    def test_nothing_can_silence_an_overlap(self):
        # Overlaps are cleared by moving or re-parenting an element. There is no
        # waiver file, so the report never offers a way to record one.
        report = self._report(self._overlapping_pair)
        self.assertIn("WARN  overlap", report)
        self.assertNotIn("waive", report)

    def test_a_parent_path_no_longer_excuses_an_overlap(self):
        # Naming used to exempt a declared child from the occlusion check. It
        # buys nothing now -- only the geometry counts.
        def build(slide):
            self._shape(slide, "chart/bar", 5.0, 0.5, 1.0, 2.0)
            self._shape(slide, "chart/bar/callout", 5.5, 0.5, 2.0, 0.5)

        self.assertIn("WARN  overlap", self._report(build))


    def test_naming_coverage_is_reported(self):
        def build(slide):
            self._shape(slide, "card", 1, 1, 2, 1)
            self._shape(slide, "Shape 9", 5, 1, 2, 1)

        report = self._report(build)
        self.assertIn("named 1 | inferred 1", report)
        self.assertIn("naming coverage 50% (1/2)", report)

    def test_missing_font_metrics_are_announced(self):
        def build(slide):
            self._text(slide, "footer/pageno", 8.95, 5.3, 0.5, 0.22,
                       "02 / 10", size=9, font="No Such Font Family", wrap=True)

        report = self._report(build)
        self.assertIn("font metrics unavailable for No Such Font Family", report)
        self.assertIn("needs 2 lines", report)   # conservative estimate still catches it

    def test_exit_code_is_zero_even_with_warnings(self):
        with tempfile.TemporaryDirectory() as temp:
            pptx = self._deck(Path(temp), self._overlapping_pair)
            with patch("sys.stdout", new=io.StringIO()) as out:
                code = layout_check.main([str(pptx)])
            self.assertEqual(code, 0)
            self.assertIn("WARN  overlap", out.getvalue())

    def test_every_check_kind_has_a_repair_hint(self):
        self.assertEqual(
            {"overlap", "clearance", "fit", "leading", "center", "baseline",
             "balance", "underfill"},
            set(layout_check.CHECK_HINTS),
        )


class SkillDocumentTests(unittest.TestCase):
    """The layout check is documented in the bundled skill, not only in code."""

    def setUp(self):
        self.skill = (ROOT / "pptx" / "SKILL.md").read_text(encoding="utf-8")

    def test_script_table_stays_one_contiguous_table(self):
        # The checker's row must not be separated from the table by a blank or
        # non-row line, which would end the table in markdown.
        lines = self.skill.splitlines()
        start = next(i for i, l in enumerate(lines) if l.startswith("| Script"))
        rows = list(itertools.takewhile(lambda l: l.startswith("|"), lines[start:]))
        self.assertTrue(any("scripts/layout_check.py" in r for r in rows))

    def test_skill_documents_the_checker_and_the_naming_convention(self):
        self.assertIn("scripts/layout_check.py", self.skill)
        self.assertIn("objectName", self.skill)
        self.assertIn("@bg/", self.skill)

if __name__ == "__main__":
    unittest.main()
